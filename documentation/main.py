import json
from pptx import Presentation
from pptx.util import Inches
import os

# Load OpenAPI JSON file
with open("api_spec.json", "r") as file:
    openapi_spec = json.load(file)

# Extract endpoints
paths = openapi_spec.get("paths", {})
schemas = openapi_spec.get("components", {}).get("schemas", {})

# Function to format JSON as a readable string
def format_json(data):
    return json.dumps(data, indent=4, ensure_ascii=False)

# Function to extract and format request body
def format_request_body(request_body):
    if not request_body:
        return "None"

    content = request_body.get("content", {})
    if "application/json" in content:
        schema = content["application/json"].get("schema", {})
        if "$ref" in schema:
            schema_name = schema["$ref"].split("/")[-1]  # Extract schema name
            formatted_schema = format_json(schemas.get(schema_name, {}))  # Pretty-print JSON
            return f"Schema: {schema_name}\n\n{formatted_schema}"
    
    return format_json(request_body)  # Fallback

# Function to extract and format response body
def format_response_body(responses):
    formatted_responses = []
    for status, response in responses.items():
        content = response.get("content", {})
        if "application/json" in content:
            schema = content["application/json"].get("schema", {})
            if "$ref" in schema:
                schema_name = schema["$ref"].split("/")[-1]
                formatted_schema = format_json(schemas.get(schema_name, {}))
                formatted_responses.append(f"Status {status}:\nSchema: {schema_name}\n\n{formatted_schema}")
            else:
                formatted_responses.append(f"Status {status}:\n{format_json(schema)}")
        else:
            formatted_responses.append(f"Status {status}: No JSON response")
    
    return "\n\n".join(formatted_responses)

# Create PowerPoint Presentation
prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]  # Title Slide Layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "API Documentation"
subtitle.text = "Generated from OpenAPI Specification"

# Loop through API paths
for endpoint, methods in paths.items():
    for method, details in methods.items():
        slide_layout = prs.slide_layouts[1]  # Title & Content Layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        # API Method & Endpoint
        title.text = f"{method.upper()} {endpoint}"
        description = details.get("description", "No description provided.")

        # Parameters
        parameters = details.get("parameters", [])
        param_text = "\n".join([f"- {p['name']} ({p['in']}): {p.get('description', 'No description')}" for p in parameters]) or "None"

        # Request & Response Body
        request_body = format_request_body(details.get("requestBody", None))
        response_body = format_response_body(details.get("responses", {}))

        # Slide Content
        content.text = (
            f"**Description:** {description}\n\n"
            f"**Parameters:**\n{param_text}\n\n"
            f"**Request Body:**\n{request_body}\n\n"
            f"**Response Body:**\n{response_body}"
        )

# Save the Presentation
output_path = "API_Documentation.pptx"
prs.save(output_path)
print(f"PowerPoint API documentation generated: {os.path.abspath(output_path)}")
